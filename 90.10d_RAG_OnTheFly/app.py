# main.py
import os
import requests
import streamlit as st
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain.vectorstores import FAISS
from langchain_text_splitters import CharacterTextSplitter
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from langchain.llms.base import LLM
from typing import Optional, List, Mapping, Any
from htmlTemplates import css, bot_template, user_template
import docx
from pptx import Presentation
from link_parser import LinkParser

# -----------------------------------
# Setup
# -----------------------------------
load_dotenv()
link_parser = LinkParser()

# -----------------------------------
# Gemini (Google Generative Language API) LLM wrapper
# -----------------------------------
class GeminiChat(LLM):
    api_key: str
    api_base: str = "https://generativelanguage.googleapis.com/v1beta"
    model: str = "models/gemini-2.5-flash"
    temperature: float = 0.0
    max_output_tokens: int = 512

    class Config:
        arbitrary_types_allowed = True

    @property
    def _llm_type(self) -> str:
        return "gemini-chat"

    @property
    def _identifying_params(self) -> Mapping[str, Any]:
        return {
            "model": self.model,
            "temperature": self.temperature,
            "max_output_tokens": self.max_output_tokens,
        }

    def _call(self, prompt: str, stop: Optional[List[str]] = None) -> str:
        # Use the API key as a query param, not a header
        url = f"{self.api_base}/{self.model}:generateContent?key={self.api_key}"

        headers = {"Content-Type": "application/json"}

        payload = {
            "contents": [{"parts": [{"text": prompt}]}],
            "generationConfig": {
                "temperature": float(self.temperature),
                "maxOutputTokens": int(self.max_output_tokens),
            },
        }

        if stop:
            payload["stopSequences"] = stop

        try:
            response = requests.post(url, headers=headers, json=payload, timeout=60)
            data = response.json()
        except Exception as e:
            return f"Error connecting to Gemini API: {e}"

        if response.status_code != 200:
            err_msg = data.get("error", {}).get("message", str(data))
            return f"Gemini API Error ({response.status_code}): {err_msg}"

        try:
            return data["candidates"][0]["content"]["parts"][0]["text"].strip()
        except Exception:
            return "No valid response from Gemini."


# -----------------------------------
# Text extraction helpers
# -----------------------------------
def getText(filename):
    doc = docx.Document(filename[0])
    return '\n'.join(p.text for p in doc.paragraphs)

def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        reader = PdfReader(pdf)
        for page in reader.pages:
            if page.extract_text():
                text += page.extract_text() + "\n"
    return text

def get_pptx_text(pptx_docs):
    prs = Presentation(pptx_docs[0])
    fullText = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                fullText.append(shape.text)
    return '\n'.join(fullText)

def get_text_from_url(url):
    try:
        return link_parser.extract_text_from_website(url)
    except Exception as e:
        st.error(f"Error extracting text: {e}")
        return None

def get_text_chunks(text):
    splitter = CharacterTextSplitter(separator="\n", chunk_size=1200, chunk_overlap=200)
    return splitter.split_text(text)

# -----------------------------------
# Vector store and conversation chain
# -----------------------------------
def get_vector_store(text_chunks):
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    return FAISS.from_texts(texts=text_chunks, embedding=embeddings)

def get_conversation_chain(vectorstore, gemini_key: str, gemini_base: Optional[str] = None):
    llm = GeminiChat(
        api_key=gemini_key,
        api_base=gemini_base or "https://generativelanguage.googleapis.com/v1beta",
        model="models/gemini-2.5-flash",
        temperature=0.2,
        max_output_tokens=1024,
    )
    memory = ConversationBufferMemory(memory_key='chat_history', return_messages=True)
    return ConversationalRetrievalChain.from_llm(llm=llm, retriever=vectorstore.as_retriever(), memory=memory)

# -----------------------------------
# Streamlit UI and handlers
# -----------------------------------
def handle_userinput(user_question):
    if st.session_state.conversation is None:
        st.warning("Please upload and analyze your document or URL first.")
        return

    result = st.session_state.conversation({"question": user_question})
    answer = result.get("answer", "")

    st.session_state.chat_history.append(("user", user_question))
    st.session_state.chat_history.append(("assistant", answer))

    for role, content in st.session_state.chat_history:
        template = user_template if role == "user" else bot_template
        st.write(template.replace("{{MSG}}", content), unsafe_allow_html=True)

def main():
    st.set_page_config(page_title="RAG (Gemini 1.5 Flash)", page_icon=":sparkles:")
    st.write(css, unsafe_allow_html=True)

    gemini_key = st.secrets.get("GEMINI_API_KEY") or os.getenv("GEMINI_API_KEY")
    gemini_base = st.secrets.get("GEMINI_API_BASE") or os.getenv("GEMINI_API_BASE")

    if not gemini_key:
        st.error("Missing GEMINI_API_KEY. Add it to Streamlit secrets or your .env file.")
        st.stop()

    if 'conversation' not in st.session_state:
        st.session_state.conversation = None
        st.session_state.chat_history = []

    st.header("RAG Chatbot (Gemini 1.5 Flash Free Tier)")
    user_input = st.text_input("Ask a question")

    if user_input:
        handle_userinput(user_input)

    with st.sidebar:
        st.subheader("Upload your documents or enter a URL")
        files = st.file_uploader("Upload PDF/DOCX/PPTX", accept_multiple_files=True)
        url = st.text_input("Enter a website URL")

        if st.button("Analyze"):
            with st.spinner("Extracting and analyzing..."):
                raw_text = None
                if files:
                    name = files[0].name.lower()
                    if name.endswith(".pdf"):
                        raw_text = get_pdf_text(files)
                    elif name.endswith(".docx"):
                        raw_text = getText(files)
                    elif name.endswith(".pptx"):
                        raw_text = get_pptx_text(files)
                    else:
                        st.error("Unsupported file format.")
                        st.stop()
                elif url:
                    raw_text = get_text_from_url(url)
                else:
                    st.error("Please upload a document or enter a URL.")
                    st.stop()

                if not raw_text:
                    st.error("No text extracted.")
                    st.stop()

                chunks = get_text_chunks(raw_text)
                store = get_vector_store(chunks)
                st.session_state.conversation = get_conversation_chain(store, gemini_key, gemini_base)

                st.success("Ready. Start chatting!")

if __name__ == "__main__":
    main()
